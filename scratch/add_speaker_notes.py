from pptx import Presentation
import os

def add_speaker_notes(pptx_path, output_path):
    prs = Presentation(pptx_path)
    
    # Speaker notes mapping (0-indexed slide numbers)
    notes_map = {
        0: (
            "Good morning. My name is Andile Nhlanhla Makuyana, registration number 25016. Today, I'll be presenting "
            "Dungeon Escape, a 2D Java game developed as a technical demonstration for ICT 2132. This project "
            "showcases the integration of core Java principles with game development mechanics."
        ),
        1: (
            "Our agenda for today covers the entire system lifecycle. We'll start with a project overview, then dive "
            "into core classes, AI algorithms, game mechanics, and data persistence. We'll conclude with a look at "
            "the ArcherBoss logic and the overall system architecture."
        ),
        2: (
            "Dungeon Escape is built entirely from first principles with no external game engine. It features a "
            "complex 16-class architecture, supporting two distinct game modes: Casual and Escape. We've "
            "implemented advanced features like A* pathfinding, procedural map generation, and a hybrid "
            "Swing-JavaFX UI system."
        ),
        3: (
            "Section 2 focuses on our core Java classes: Main, GameMenu, KeyHandler, and Player. These classes "
            "form the backbone of the game's startup, menu interface, and user control systems."
        ),
        4: (
            "In Main.java, we use a standard JFrame bootstrap to initialize the Swing window. A unique feature here "
            "is the 'Cross-Thread Bridge' in quitToMenu(), which uses Platform.runLater to safely transition from "
            "the Swing game loop to the JavaFX UI thread. For the Menu, we use JavaFX property bindings to ensure "
            "the UI scales correctly while maintaining a 4:3 aspect ratio."
        ),
        5: (
            "KeyHandler.java uses a polling model with boolean flags to handle multi-key input smoothly. For the "
            "Player, we've implemented axis-independent collision. As you can see in the code snippet, we test X "
            "and Y movement separately. This is crucial for 'wall-sliding'—if a player moves diagonally against "
            "a wall, they will still slide along it rather than coming to a full stop."
        ),
        6: (
            "Section 3 covers our AI and procedural generation logic, focusing on how we manage complexity in "
            "the game world and NPC behavior."
        ),
        7: (
            "The Enemy AI uses a custom A* pathfinding implementation. We select the node with the lowest fCost "
            "each frame using a Manhattan distance heuristic. The AI is a finite state machine with two modes: "
            "CHASE and WANDER. To optimize performance, we recalculate the path every 15 frames rather than "
            "every single frame, and we've implemented a repulsion logic to keep enemies from overlapping."
        ),
        8: (
            "The Node class is a clean POJO that forms a linked-list path chain. It uses direct field access for "
            "performance and adheres to the Open/Closed principle for its cost calculations. For the Map, we use "
            "procedural generation where the 'wall chance' scales with the level. We've also implemented a 'Drunk-walk' "
            "path carver to ensure there is always a traversable 2x2 corridor from the spawn to the exit."
        ),
        10: (
            "Section 4 details our game mechanics, specifically our projectile system and interactive collectibles."
        ),
        11: (
            "Arrow.java uses vector movement derived from Math.atan2. It features dual constructors—one for direct "
            "targeting and one for pre-computed spread shots. The ExplosionTracker is a 4-phase state machine that "
            "lerps from a red 'tracking' state to a green 'locked' state before detonating, using AABB blast rectangles "
            "for uniform collision detection across all entity types."
        ),
        13: (
            "The Crystal class demonstrates the Single Responsibility Principle. It manages its own boolean lifecycle "
            "and rendering. To make pickup feel more responsive, we've implemented 'Expanded Hit Bounds' where the "
            "collision box is 18 pixels larger than the visual sprite on each side, avoiding frustration for the player."
        ),
        14: (
            "Section 5 covers Data and Persistence, specifically how we handle score saving and efficient map rendering."
        ),
        15: (
            "ScoreEntry implements Serializable and uses a custom CSV encoding for storage. ScoreManager handles "
            "file I/O with defensive parsing to remain backward compatible with older save files. For Tile.java, "
            "we use the Flyweight pattern—instead of 300 unique objects for a 20x15 grid, we use just 6 shared "
            "tile instances referenced by an integer array, drastically reducing the memory footprint."
        ),
        16: (
            "Section 6 brings everything together with the ArcherBoss AI and the central GamePanel orchestrator."
        ),
        17: (
            "The ArcherBoss has three distinct attack phases: a Shotgun spread, a high-intensity 'Frenzy' sprinkler "
            "triggered by crystal hits, and a Machine-Gun burst. You can see the code for the Frenzy phase here, "
            "where it spawns arrows at 90-degree offsets while rotating. This creates a challenging 'bullet-hell' "
            "experience for the player."
        ),
        18: (
            "Advanced boss mechanics include A* movement for retreating, random taunts, and twin-boss separation "
            "logic for Level 10. We also use visual effects like white damage flashes and alpha-faded ghost trails "
            "to provide immediate feedback to the player during intense combat."
        ),
        19: (
            "GamePanel.java is the spine of the architecture. It owns every collection and drives all timers, "
            "including the slowTimer for power-ups and the ghostFreezeTimer. It manages the GameMode enum to "
            "seamlessly switch between Casual and Escape mode logic."
        ),
        20: (
            "This table provides a component overview of the system architecture. It highlights how each class has "
            " a specific responsibility, from the JavaFX-based GameMenu to the A* pathfinding in Enemy.java and "
            "the CSV serialization in ScoreManager."
        ),
        21: (
            "The system workflow starts with the JavaFX menu and leads into procedural map generation. Players explore, "
            "collect crystals, and battle bosses before saving their score. We've also included several debug "
            "key-binds like 'G' for God Mode and 'L' for skipping to the final chase, which were vital for testing."
        ),
        22: (
            "Reflecting on the project, our modular SRP architecture and hand-implemented A* are major strengths. "
            "For future improvements, we're looking at moving from CSV to SQLite for score persistence and "
            "implementing a centralized AssetManager to further optimize how we handle images and resources."
        ),
        23: (
            "In conclusion, Dungeon Escape proves that a feature-complete game can be built from scratch in Java "
            "using clean OOP architecture and hand-rolled algorithms. The 16-class codebase covers 11 chapters "
            "of Java theory and 4 design patterns. Thank you for listening, and I'm happy to take any questions."
        )
    }
    
    for slide_idx, notes in notes_map.items():
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            if text_frame is not None:
                text_frame.text = notes
            else:
                # Fallback: find the body placeholder in the notes slide
                for shape in notes_slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.type == 2: # 2 is BODY
                        shape.text = notes
                        break
            print(f"Added notes to Slide {slide_idx + 1}")
    
    prs.save(output_path)
    print(f"Saved modified presentation to: {output_path}")

if __name__ == "__main__":
    input_file = "Dungeon_Escape_Presentation_v2.pptx"
    output_file = "Dungeon_Escape_Presentation_v2_with_Notes.pptx"
    add_speaker_notes(input_file, output_file)
